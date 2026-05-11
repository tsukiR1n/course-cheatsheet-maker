"""Extract course materials into Markdown text files for agent analysis.

This script is intentionally conservative: it records every source file it
finds, documents extraction failures, and never asks downstream agents to infer
document content from filenames alone.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import subprocess
import sys
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path


SUPPORTED_TEXT = {".md", ".txt"}
SUPPORTED_BINARY = {".pdf", ".pptx", ".docx"}
SUPPORTED_EXTENSIONS = SUPPORTED_TEXT | SUPPORTED_BINARY
SKILL_ROOT = Path(__file__).resolve().parents[1]
REQUIREMENTS_PATH = SKILL_ROOT / "requirements.txt"
NO_AUTO_INSTALL_ENV = "COURSE_CHEATSHEET_NO_AUTO_INSTALL"
INSTALL_COMMAND = f"{sys.executable} -m pip install -r {REQUIREMENTS_PATH}"
NO_TEXT_MESSAGE = "No selectable text was extracted. OCR is out of scope for v1."


QUESTION_TYPES = [
    ("past_paper", re.compile(r"\bpast[_ -]?paper\b|(?<![a-z])exam[_ -]?\d{4}", re.IGNORECASE)),
    (
        "final_exam",
        re.compile(
            r"\bfinal[_ -]?exam\b|\bthis examination takes\b|\bfor examiner's use only\b|\bend of paper\b",
            re.IGNORECASE,
        ),
    ),
    ("sample_exam", re.compile(r"\bsample[_ -]?exam\b", re.IGNORECASE)),
    ("mock_exam", re.compile(r"\bmock[_ -]?exam\b", re.IGNORECASE)),
    ("practice_exam", re.compile(r"\bpractice[_ -]?exam\b", re.IGNORECASE)),
    ("quiz", re.compile(r"\b(quiz|weekly[_ -]?quiz)\b", re.IGNORECASE)),
    ("assignment", re.compile(r"\bassignment\b", re.IGNORECASE)),
    ("homework", re.compile(r"\b(homework|hw)\b", re.IGNORECASE)),
    ("problem_set", re.compile(r"\bproblem[_ -]?set\b", re.IGNORECASE)),
    ("exercise_sheet", re.compile(r"\bexercise[_ -]?sheet\b", re.IGNORECASE)),
    ("workshop", re.compile(r"\bworkshop\b", re.IGNORECASE)),
    ("tutorial", re.compile(r"\btutorial\b", re.IGNORECASE)),
]


@dataclass
class ExtractionRecord:
    source_id: str
    source_path: Path
    source_folder: str
    relative_path: str
    output_path: Path
    inferred_source_type: str
    status: str
    extractor: str
    extracted_units: str
    message: str
    text: str = ""


@dataclass
class DependencyBootstrapResult:
    status: str
    message: str


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Extract course materials into working/extracted Markdown files."
    )
    parser.add_argument("course_name", help="Course folder name under --root.")
    parser.add_argument("--root", default="courses", help="Courses root directory.")
    return parser.parse_args()


def json_scalar(value: str) -> str:
    return json.dumps(value, ensure_ascii=False)


def safe_output_name(relative_path: str) -> str:
    path = Path(relative_path)
    stem = str(path.with_suffix("")).replace("\\", "__").replace("/", "__")
    stem = re.sub(r"[^A-Za-z0-9._-]+", "_", stem).strip("._")
    if not stem:
        stem = "material"
    digest = hashlib.sha1(relative_path.encode("utf-8")).hexdigest()[:10]
    return f"{stem}__{digest}.md"


def read_plain_text(path: Path) -> tuple[str, str]:
    for encoding in ["utf-8-sig", "utf-8", "cp1252"]:
        try:
            return path.read_text(encoding=encoding), f"Read as {encoding}."
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace"), (
        "Read as utf-8 with replacement characters after encoding detection failed."
    )


def chunk_text(text: str, max_chars: int = 12000) -> tuple[str, int]:
    if len(text) <= max_chars:
        return text, 1 if text else 0

    chunks: list[str] = []
    start = 0
    chunk_number = 1
    while start < len(text):
        end = min(start + max_chars, len(text))
        if end < len(text):
            next_break = text.rfind("\n\n", start, end)
            if next_break > start + max_chars // 2:
                end = next_break
        body = text[start:end].strip()
        if body:
            chunks.append(f"## Chunk {chunk_number}\n\n{body}")
            chunk_number += 1
        start = end
    return "\n\n".join(chunks), len(chunks)


def ensure_extraction_dependencies() -> DependencyBootstrapResult:
    script = Path(__file__).resolve().parent / "ensure_dependencies.py"
    if os.environ.get(NO_AUTO_INSTALL_ENV) == "1":
        message = (
            f"Automatic dependency installation skipped because {NO_AUTO_INSTALL_ENV}=1."
        )
        print(message)
        return DependencyBootstrapResult(status="skipped", message=message)

    print("Checking extraction dependencies before material extraction...")
    result = subprocess.run([sys.executable, str(script)], check=False)
    if result.returncode == 0:
        return DependencyBootstrapResult(
            status="success",
            message="Dependency check completed; available dependencies were confirmed or installed.",
        )

    return DependencyBootstrapResult(
        status="failed",
        message=(
            "Dependency bootstrap failed. Extraction will continue and document "
            "per-file missing-dependency failures where applicable. "
            f"Install manually with: {INSTALL_COMMAND}"
        ),
    )


def extract_pdf(path: Path) -> tuple[str, str, str, str]:
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError:
        return (
            "",
            "failed",
            "0 pages",
            f"Missing optional dependency PyMuPDF. Install with: {INSTALL_COMMAND}",
        )

    try:
        chunks: list[str] = []
        text_pages = 0
        with fitz.open(path) as document:
            page_count = len(document)
            for index, page in enumerate(document, start=1):
                text = page.get_text("text").strip()
                if text:
                    text_pages += 1
                    chunks.append(f"## Page {index}\n\n{text}")
                else:
                    chunks.append(
                        f"## Page {index}\n\n<!-- No selectable text on this page. -->"
                    )
        extracted = "\n\n".join(chunks).strip()
        units = f"{page_count} page{'s' if page_count != 1 else ''}"
        if text_pages == 0:
            return extracted, "no_selectable_text", units, NO_TEXT_MESSAGE
        return (
            extracted,
            "success",
            units,
            f"Extracted selectable text from {text_pages} of {page_count} page(s) with PyMuPDF.",
        )
    except Exception as exc:  # pragma: no cover - depends on third-party parser internals.
        return "", "failed", "0 pages", f"PyMuPDF extraction failed: {exc}"


def extract_pptx(path: Path) -> tuple[str, str, str, str]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        return (
            "",
            "failed",
            "0 slides",
            f"Missing optional dependency python-pptx. Install with: {INSTALL_COMMAND}",
        )

    try:
        presentation = Presentation(path)
        slides: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            lines.append(" | ".join(cells))
            if lines:
                slides.append(f"## Slide {slide_index}\n\n" + "\n\n".join(lines))
            else:
                slides.append(
                    f"## Slide {slide_index}\n\n<!-- No visible text on this slide. -->"
                )
        extracted = "\n\n".join(slides).strip()
        slide_count = len(presentation.slides)
        text_slides = sum(
            1 for slide in slides if "<!-- No visible text on this slide. -->" not in slide
        )
        units = f"{slide_count} slide{'s' if slide_count != 1 else ''}"
        if text_slides == 0:
            return extracted, "no_selectable_text", units, "No slide text was extracted."
        return (
            extracted,
            "success",
            units,
            f"Extracted visible text from {text_slides} of {slide_count} slide(s). "
            "Speaker notes may not be fully extracted in v1.",
        )
    except Exception as exc:  # pragma: no cover - depends on third-party parser internals.
        return "", "failed", "0 slides", f"python-pptx extraction failed: {exc}"


def extract_docx(path: Path) -> tuple[str, str, str, str]:
    try:
        from docx import Document  # type: ignore[import-not-found]
    except ImportError:
        return (
            "",
            "failed",
            "0 sections",
            f"Missing optional dependency python-docx. Install with: {INSTALL_COMMAND}",
        )

    try:
        document = Document(path)
        lines: list[str] = ["## Document Text"]
        paragraph_count = 0
        table_count = 0
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = paragraph.style.name if paragraph.style else ""
            if style_name.lower().startswith("heading"):
                lines.append(f"### Heading: {text}")
            else:
                lines.append(f"Paragraph: {text}")
            paragraph_count += 1
        for table in document.tables:
            table_lines: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    table_lines.append(" | ".join(cells))
            if table_lines:
                table_count += 1
                lines.append("")
                lines.append(f"Table: {table_count}")
                lines.extend(table_lines)
        extracted = "\n\n".join(lines).strip()
        units = (
            f"{paragraph_count} paragraph{'s' if paragraph_count != 1 else ''}, "
            f"{table_count} table{'s' if table_count != 1 else ''}"
        )
        if paragraph_count == 0 and table_count == 0:
            return "", "no_selectable_text", units, "No document text was extracted."
        return (
            extracted,
            "success",
            units,
            "Extracted structured paragraphs, headings, and table text with python-docx.",
        )
    except Exception as exc:  # pragma: no cover - depends on third-party parser internals.
        return "", "failed", "0 sections", f"python-docx extraction failed: {exc}"


def infer_question_type(relative_path: str, extracted_text: str = "") -> str:
    signal = f"{relative_path}\n{extracted_text[:4000]}"
    for source_type, pattern in QUESTION_TYPES:
        if pattern.search(signal):
            return source_type
    return "unknown_question_source"


def extract_file(
    source_id: str,
    source_path: Path,
    source_folder: str,
    relative_path: str,
    output_path: Path,
) -> ExtractionRecord:
    suffix = source_path.suffix.lower()
    inferred_source_type = "knowledge_source"
    text = ""
    status = "unsupported"
    extractor = "none"
    extracted_units = "0 units"
    message = f"Unsupported file extension: {suffix or '(none)'}."

    if suffix in SUPPORTED_TEXT:
        extractor = "stdlib"
        try:
            text, message = read_plain_text(source_path)
            text, chunk_count = chunk_text(text)
            extracted_units = f"{chunk_count} chunk{'s' if chunk_count != 1 else ''}"
            status = "success" if text.strip() else "no_selectable_text"
            if status == "no_selectable_text":
                message = "Text file is empty."
        except OSError as exc:
            status = "failed"
            message = f"Could not read text file: {exc}"
    elif suffix == ".pdf":
        extractor = "PyMuPDF"
        text, status, extracted_units, message = extract_pdf(source_path)
    elif suffix == ".pptx":
        extractor = "python-pptx"
        text, status, extracted_units, message = extract_pptx(source_path)
    elif suffix == ".docx":
        extractor = "python-docx"
        text, status, extracted_units, message = extract_docx(source_path)

    if source_folder == "questions":
        inferred_source_type = infer_question_type(relative_path, text if status == "success" else "")

    return ExtractionRecord(
        source_id=source_id,
        source_path=source_path,
        source_folder=source_folder,
        relative_path=relative_path,
        output_path=output_path,
        inferred_source_type=inferred_source_type,
        status=status,
        extractor=extractor,
        extracted_units=extracted_units,
        message=message,
        text=text if status in {"success", "no_selectable_text"} else "",
    )


def write_extracted_file(record: ExtractionRecord) -> None:
    record.output_path.parent.mkdir(parents=True, exist_ok=True)
    metadata = {
        "source_id": record.source_id,
        "original_file": record.source_path.name,
        "source_folder": record.source_folder,
        "relative_path": record.relative_path,
        "inferred_source_type": record.inferred_source_type,
        "status": record.status,
        "extractor": record.extractor,
        "extracted_units": record.extracted_units,
        "message": record.message,
    }
    lines = ["---"]
    lines.extend(f"{key}: {json_scalar(value)}" for key, value in metadata.items())
    lines.append("---")
    lines.append("")
    lines.append(f"# {record.source_path.name}")
    lines.append("")
    if record.text:
        lines.append("## Extracted Text")
        lines.append("")
        lines.append(record.text)
    else:
        lines.append("## Extraction Note")
        lines.append("")
        lines.append(record.message)
    record.output_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def material_files(material_dir: Path) -> list[Path]:
    if not material_dir.exists():
        return []
    return sorted(path for path in material_dir.rglob("*") if path.is_file())


def write_report(
    course_name: str,
    report_path: Path,
    records: list[ExtractionRecord],
    dependency_bootstrap: DependencyBootstrapResult,
) -> None:
    counts: dict[str, int] = {}
    for record in records:
        counts[record.status] = counts.get(record.status, 0) + 1

    lines = [
        f"# Extraction Report: {course_name}",
        "",
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        "",
        "## Dependency Bootstrap",
        "",
        f"- status: {dependency_bootstrap.status}",
        f"- message: {dependency_bootstrap.message}",
        "",
        "Optional dependency install command:",
        "",
        "```bash",
        INSTALL_COMMAND,
        "```",
        "",
        "OCR is out of scope for v1. Scanned PDFs may be reported as `no_selectable_text`.",
        "",
        "## Summary",
        "",
    ]
    if records:
        for status in ["success", "failed", "unsupported", "no_selectable_text"]:
            lines.append(f"- {status}: {counts.get(status, 0)}")
    else:
        lines.append("- No material files found.")

    lines.extend(
        [
            "",
            "## Files",
            "",
            "| Source ID | File | Folder | Inferred type | Extracted units | Status | Output | Message |",
            "| --- | --- | --- | --- | --- | --- | --- | --- |",
        ]
    )
    for record in records:
        output = record.output_path.as_posix()
        message = record.message.replace("|", "\\|").replace("\n", " ")
        lines.append(
            "| "
            f"{record.source_id} | `{record.relative_path}` | {record.source_folder} | "
            f"{record.inferred_source_type} | {record.extracted_units} | {record.status} | "
            f"`{output}` | {message} |"
        )

    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def run(course_name: str, root: Path) -> int:
    course_root = root / course_name
    materials_root = course_root / "materials"
    working_root = course_root / "working"
    extracted_root = working_root / "extracted"

    if not course_root.exists():
        print(f"ERROR: Course folder not found: {course_root}")
        return 2
    if not materials_root.exists():
        print(f"ERROR: Materials folder not found: {materials_root}")
        return 2

    dependency_bootstrap = ensure_extraction_dependencies()

    records: list[ExtractionRecord] = []
    for source_folder in ["knowledge", "questions"]:
        source_dir = materials_root / source_folder
        output_dir = extracted_root / source_folder
        output_dir.mkdir(parents=True, exist_ok=True)

        prefix = "K" if source_folder == "knowledge" else "Q"
        for index, source_path in enumerate(material_files(source_dir), start=1):
            source_id = f"{prefix}{index:02d}"
            relative_path = source_path.relative_to(source_dir).as_posix()
            output_path = output_dir / safe_output_name(relative_path)
            record = extract_file(
                source_id, source_path, source_folder, relative_path, output_path
            )
            write_extracted_file(record)
            records.append(record)

    report_path = working_root / "extraction_report.md"
    write_report(course_name, report_path, records, dependency_bootstrap)

    print(f"Wrote {report_path}")
    print(f"Processed {len(records)} material file(s).")
    failures = [record for record in records if record.status in {"failed", "no_selectable_text"}]
    if failures:
        print(
            f"Documented {len(failures)} extraction issue(s). "
            "Review extraction_report.md before topic analysis."
        )
    return 0


def main() -> int:
    args = parse_args()
    return run(args.course_name, Path(args.root))


if __name__ == "__main__":
    sys.exit(main())
